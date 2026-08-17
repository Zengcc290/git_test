"""Format-aware parsing into DocumentBlock objects."""

from __future__ import annotations

import ast
import hashlib
import io
import logging
import re
import tokenize
from collections.abc import Iterable, Iterator
from dataclasses import replace
from pathlib import Path

from .cleaning import clean_text
from .extractors import iter_document_text
from .json_parser import (
    DEFAULT_JSON_RECORD_PROBE_SIZE,
    DEFAULT_MAX_JSON_SIZE,
    JsonProfile,
    iter_json_record_text,
)
from .models import DocumentBlock, ExtractedDocument
from .embedding import build_document_embedding_input


logger = logging.getLogger(__name__)

CPP_LANGUAGE = {
    "c": "C",
    "h": "C/C++",
    "cc": "C++",
    "cpp": "C++",
    "cxx": "C++",
    "hpp": "C++",
}


class _BlockBuilder:
    """Create stable, unique block ids in source order."""

    def __init__(self, path: Path) -> None:
        self.path = path
        self.sequence = 0

    def make(
        self,
        block_type: str,
        content: str,
        *,
        language: str | None = None,
        heading_path: tuple[str, ...] = (),
        symbol_path: tuple[str, ...] = (),
        start_line: int | None = None,
        end_line: int | None = None,
        page_number: int | None = None,
        hard_boundary_before: bool = False,
        hard_boundary_after: bool = False,
        record_path: str | None = None,
        slide_number: int | None = None,
        shape_index: int | None = None,
        module_name: str | None = None,
        parameters: tuple[str, ...] = (),
        docstring: str | None = None,
        comments: tuple[str, ...] = (),
        parser: str = "",
    ) -> DocumentBlock:
        identity = (
            f"{self.path}\0{self.sequence}\0{block_type}\0{start_line}\0"
            f"{page_number}\0{slide_number}\0{record_path or ''}"
        )
        block_id = hashlib.sha256(identity.encode("utf-8")).hexdigest()[:24]
        self.sequence += 1
        return DocumentBlock(
            block_id=block_id,
            path=str(self.path),
            block_type=block_type,
            language=language,
            heading_path=heading_path,
            symbol_path=symbol_path,
            content=content,
            start_line=start_line,
            end_line=end_line,
            page_number=page_number,
            hard_boundary_before=hard_boundary_before,
            hard_boundary_after=hard_boundary_after,
            record_path=record_path,
            slide_number=slide_number,
            shape_index=shape_index,
            module_name=module_name,
            parameters=parameters,
            docstring=docstring,
            comments=comments,
            parser=parser,
        )


def _numbered_lines(parts: Iterable[str]) -> list[tuple[int, str]]:
    text = "".join(parts).replace("\r\n", "\n").replace("\r", "\n")
    return list(enumerate(text.splitlines(), start=1))


_MD_HEADING = re.compile(r"^(#{1,3})\s+(.+?)\s*#*\s*$")
_MD_LIST = re.compile(r"^\s*(?:[-+*]|\d+[.)])\s+")
_MD_TABLE_RULE = re.compile(
    r"^\s*\|?\s*:?-{3,}:?\s*(?:\|\s*:?-{3,}:?\s*)+\|?\s*$"
)


def _markdown_node_kind(lines: list[tuple[int, str]], index: int) -> str:
    line = lines[index][1]
    stripped = line.lstrip()
    if stripped.startswith(("```", "~~~")):
        return "code"
    if _MD_LIST.match(line):
        return "list"
    if stripped.startswith(">"):
        return "quote"
    if (
        "|" in line
        and index + 1 < len(lines)
        and _MD_TABLE_RULE.match(lines[index + 1][1])
    ):
        return "table"
    return "paragraph"


def _take_markdown_node(
    lines: list[tuple[int, str]], index: int
) -> tuple[str, str, int, int, int]:
    """Return kind, content, first line, last line and next cursor."""

    start_line, first = lines[index]
    kind = _markdown_node_kind(lines, index)
    collected = [first]
    cursor = index + 1

    if kind == "code":
        fence = first.lstrip()[:3]
        while cursor < len(lines):
            collected.append(lines[cursor][1])
            cursor += 1
            if collected[-1].lstrip().startswith(fence):
                break
    elif kind == "table":
        while cursor < len(lines) and "|" in lines[cursor][1]:
            collected.append(lines[cursor][1])
            cursor += 1
    else:
        while cursor < len(lines):
            current = lines[cursor][1]
            if not current.strip() or _MD_HEADING.match(current):
                break
            current_kind = _markdown_node_kind(lines, cursor)
            continuation = current.startswith((" ", "\t"))
            if current_kind != kind and not (kind in {"list", "quote"} and continuation):
                break
            collected.append(current)
            cursor += 1

    return kind, "\n".join(collected).strip(), start_line, lines[cursor - 1][0], cursor


def _parse_markdown(document: ExtractedDocument) -> Iterator[DocumentBlock]:
    builder = _BlockBuilder(document.path)
    lines = _numbered_lines(iter_document_text(document))
    headings: list[str] = []
    unresolved_heading: tuple[int, tuple[str, ...]] | None = None
    current_section_has_body = False
    cursor = 0

    while cursor < len(lines):
        line_number, line = lines[cursor]
        heading = _MD_HEADING.match(line)
        if heading:
            if unresolved_heading and not current_section_has_body:
                previous_line, previous_path = unresolved_heading
                yield builder.make(
                    "heading",
                    "",
                    heading_path=previous_path,
                    start_line=previous_line,
                    end_line=previous_line,
                    hard_boundary_before=True,
                    hard_boundary_after=True,
                    parser="markdown-block",
                )
            level = len(heading.group(1))
            headings = headings[: level - 1]
            headings.append(heading.group(2).strip())
            unresolved_heading = (line_number, tuple(headings))
            current_section_has_body = False
            cursor += 1
            continue
        if not line.strip():
            cursor += 1
            continue

        kind, content, start, end, cursor = _take_markdown_node(lines, cursor)
        current_section_has_body = True
        hard = kind in {"code", "table"}
        yield builder.make(
            kind,
            content,
            language="Markdown" if kind == "code" else None,
            heading_path=tuple(headings),
            start_line=start,
            end_line=end,
            hard_boundary_before=hard,
            hard_boundary_after=hard,
            parser="markdown-block",
        )

    if unresolved_heading and not current_section_has_body:
        line_number, heading_path = unresolved_heading
        yield builder.make(
            "heading",
            "",
            heading_path=heading_path,
            start_line=line_number,
            end_line=line_number,
            hard_boundary_before=True,
            hard_boundary_after=True,
            parser="markdown-block",
        )


def _parse_text(document: ExtractedDocument) -> Iterator[DocumentBlock]:
    builder = _BlockBuilder(document.path)
    lines = _numbered_lines(iter_document_text(document))
    cursor = 0
    while cursor < len(lines):
        while cursor < len(lines) and not lines[cursor][1].strip():
            cursor += 1
        if cursor >= len(lines):
            break
        start = lines[cursor][0]
        body: list[str] = []
        while cursor < len(lines) and lines[cursor][1].strip():
            body.append(lines[cursor][1])
            cursor += 1
        content = clean_text("\n".join(body))
        if content:
            yield builder.make(
                "paragraph",
                content,
                start_line=start,
                end_line=start + len(body) - 1,
                parser="text-line",
            )


def _parse_json(
    document: ExtractedDocument,
    profile: JsonProfile,
    *,
    max_json_size: int,
    record_probe_size: int,
) -> Iterator[DocumentBlock]:
    builder = _BlockBuilder(document.path)
    record_parts: list[str] = []
    locator: str | None = None

    def emit_record() -> DocumentBlock | None:
        content = clean_text("".join(record_parts))
        if not content:
            return None
        return builder.make(
            "json-record",
            content,
            record_path=locator,
            hard_boundary_before=True,
            hard_boundary_after=True,
            parser="json-stream",
        )

    for part in iter_json_record_text(
        document.path,
        profile,
        max_size=max_json_size,
        record_probe_size=record_probe_size,
    ):
        if part.record_start:
            if record_parts:
                block = emit_record()
                if block:
                    yield block
            record_parts = []
            locator = part.record_path
        record_parts.append(str(part))
        if part.record_end:
            block = emit_record()
            if block:
                yield block
            record_parts = []
            locator = None
    if record_parts:
        block = emit_record()
        if block:
            yield block


def _parse_pdf(document: ExtractedDocument) -> Iterator[DocumentBlock]:
    try:
        from pypdf import PdfReader
    except ImportError as exc:  # pragma: no cover
        raise RuntimeError("PDF 抽取需要 pypdf，请先安装 requirements.txt") from exc

    builder = _BlockBuilder(document.path)
    reader = PdfReader(str(document.path))
    for page_number, page in enumerate(reader.pages, start=1):
        try:
            page_text = page.extract_text() or ""
        except Exception:
            logger.exception("抽取 PDF %s 的第 %s 页失败", document.path, page_number)
            continue
        paragraphs = [
            clean_text(value)
            for value in re.split(r"\n\s*\n", page_text)
            if value.strip()
        ]
        for position, paragraph in enumerate(paragraphs):
            yield builder.make(
                "paragraph",
                paragraph,
                page_number=page_number,
                hard_boundary_before=position == 0,
                hard_boundary_after=position == len(paragraphs) - 1,
                parser="pypdf",
            )


def _shape_text(shape, group_type) -> tuple[str, str]:
    values: list[str] = []
    kind = "text-box"
    if getattr(shape, "has_text_frame", False) and shape.text.strip():
        values.append(shape.text.strip())
    if getattr(shape, "has_table", False):
        kind = "table"
        for row in shape.table.rows:
            values.append("\t".join(cell.text.strip() for cell in row.cells))
    if getattr(shape, "shape_type", None) == group_type:
        kind = "group"
        for child in shape.shapes:
            child_text, _ = _shape_text(child, group_type)
            if child_text:
                values.append(child_text)
    return clean_text("\n".join(values)), kind


def _parse_pptx(document: ExtractedDocument) -> Iterator[DocumentBlock]:
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError as exc:  # pragma: no cover
        raise RuntimeError("PPTX 抽取需要 python-pptx，请先安装 requirements.txt") from exc

    builder = _BlockBuilder(document.path)
    presentation = Presentation(str(document.path))
    for slide_number, slide in enumerate(presentation.slides, start=1):
        slide_blocks: list[DocumentBlock] = []
        title_shape = slide.shapes.title
        title_shape_id = title_shape.shape_id if title_shape is not None else None
        for shape_index, shape in enumerate(slide.shapes):
            content, kind = _shape_text(shape, MSO_SHAPE_TYPE.GROUP)
            if not content:
                continue
            if title_shape_id is not None and shape.shape_id == title_shape_id:
                kind = "title"
            slide_blocks.append(
                builder.make(
                    kind,
                    content,
                    slide_number=slide_number,
                    shape_index=shape_index,
                    parser="python-pptx",
                )
            )
        try:
            notes = clean_text(slide.notes_slide.notes_text_frame.text)
        except (AttributeError, ValueError):
            notes = ""
        if notes:
            slide_blocks.append(
                builder.make(
                    "notes",
                    notes,
                    slide_number=slide_number,
                    shape_index=-1,
                    parser="python-pptx",
                )
            )
        for position, block in enumerate(slide_blocks):
            yield replace(
                block,
                hard_boundary_before=position == 0,
                hard_boundary_after=position == len(slide_blocks) - 1,
            )


def _source_text(document: ExtractedDocument) -> str:
    if document.text is not None:
        return document.text
    if document.file_type == "py":
        try:
            with tokenize.open(document.path) as source:
                return source.read()
        except (SyntaxError, UnicodeDecodeError):
            pass
    return "".join(iter_document_text(document))


def _source_lines(source: str, start: int, end: int) -> str:
    return "".join(source.splitlines(keepends=True)[start - 1 : end]).rstrip()


def _source_comments(source: str, start: int, end: int) -> tuple[str, ...]:
    found: list[str] = []
    try:
        tokens = tokenize.generate_tokens(io.StringIO(source).readline)
        for token in tokens:
            if token.type == tokenize.COMMENT and start <= token.start[0] <= end:
                found.append(token.string)
    except (IndentationError, tokenize.TokenError):
        pass
    return tuple(found)


def _arguments(node: ast.FunctionDef | ast.AsyncFunctionDef) -> tuple[str, ...]:
    result = [item.arg for item in (*node.args.posonlyargs, *node.args.args)]
    if node.args.vararg:
        result.append(f"*{node.args.vararg.arg}")
    elif node.args.kwonlyargs:
        result.append("*")
    result.extend(item.arg for item in node.args.kwonlyargs)
    if node.args.kwarg:
        result.append(f"**{node.args.kwarg.arg}")
    return tuple(result)


def _python_items(
    statements: list[ast.stmt], parents: tuple[str, ...] = ()
) -> Iterator[tuple[ast.AST, str, tuple[str, ...]]]:
    for statement in statements:
        if isinstance(statement, (ast.Import, ast.ImportFrom)):
            yield statement, "imports", parents
        elif isinstance(statement, (ast.Assign, ast.AnnAssign)) and not parents:
            yield statement, "constant", ()
        elif isinstance(statement, ast.ClassDef):
            path = (*parents, statement.name)
            yield statement, "class", path
            yield from _python_items(statement.body, path)
        elif isinstance(statement, (ast.FunctionDef, ast.AsyncFunctionDef)):
            path = (*parents, statement.name)
            yield statement, "method" if parents else "function", path


def _fallback_lines(
    document: ExtractedDocument,
    source: str,
    *,
    language: str,
    builder: _BlockBuilder,
    parser: str = "fallback-line",
) -> Iterator[DocumentBlock]:
    lines = source.splitlines(keepends=True)
    for start_index in range(0, len(lines), 80):
        content = "".join(lines[start_index : start_index + 80]).rstrip()
        if content.strip():
            yield builder.make(
                "fallback-line",
                content,
                language=language,
                start_line=start_index + 1,
                end_line=min(start_index + 80, len(lines)),
                hard_boundary_before=True,
                hard_boundary_after=True,
                module_name=document.path.stem if language == "Python" else None,
                parser=parser,
            )


def _parse_python(document: ExtractedDocument) -> Iterator[DocumentBlock]:
    builder = _BlockBuilder(document.path)
    source = _source_text(document)
    try:
        module = ast.parse(source, filename=str(document.path), type_comments=True)
    except SyntaxError as exc:
        logger.warning(
            "Python AST 解析失败，按行降级：%s:%s：%s",
            document.path,
            exc.lineno or "?",
            exc.msg,
        )
        yield from _fallback_lines(
            document, source, language="Python", builder=builder
        )
        return

    module_docstring = ast.get_docstring(module, clean=False)
    if module_docstring and module.body:
        expression = module.body[0]
        yield builder.make(
            "module",
            ast.get_source_segment(source, expression) or module_docstring,
            language="Python",
            symbol_path=(document.path.stem,),
            start_line=expression.lineno,
            end_line=getattr(expression, "end_lineno", expression.lineno),
            module_name=document.path.stem,
            docstring=module_docstring,
            parser="python-ast",
        )

    for node, kind, symbol_path in _python_items(module.body):
        decorators = getattr(node, "decorator_list", ())
        start = min((item.lineno for item in decorators), default=node.lineno)
        end = getattr(node, "end_lineno", node.lineno)
        raw_source = _source_lines(source, start, end)
        if not raw_source:
            continue
        function = node if isinstance(node, (ast.FunctionDef, ast.AsyncFunctionDef)) else None
        yield builder.make(
            kind,
            raw_source,
            language="Python",
            symbol_path=symbol_path,
            start_line=start,
            end_line=end,
            hard_boundary_before=True,
            hard_boundary_after=True,
            module_name=document.path.stem,
            parameters=_arguments(function) if function else (),
            docstring=ast.get_docstring(node, clean=False)
            if isinstance(node, (ast.ClassDef, ast.FunctionDef, ast.AsyncFunctionDef))
            else None,
            comments=_source_comments(source, start, end),
            parser="python-ast",
        )


def _brace_change(line: str) -> int:
    without_literals = re.sub(
        r'"(?:\\.|[^"\\])*"|\'(?:\\.|[^\'\\])*\'', "", line
    )
    without_comment = without_literals.split("//", 1)[0]
    return without_comment.count("{") - without_comment.count("}")


def _closing_line(lines: list[str], start: int) -> int:
    depth = 0
    opened = False
    for position in range(start, len(lines)):
        code = lines[position].split("//", 1)[0]
        opened = opened or "{" in code
        depth += _brace_change(lines[position])
        if opened and depth <= 0:
            return position
    return len(lines) - 1


_CPP_CONTAINER = re.compile(
    r"\b(namespace|class|struct|enum)\s+([A-Za-z_]\w*(?:::\w+)*)[^;{]*\{"
)
_CPP_FUNCTION = re.compile(
    r"(?:^|\s)([~A-Za-z_]\w*(?:::\w+)*)\s*\([^;{}]*\)\s*"
    r"(?:const\s*)?(?:noexcept(?:\([^)]*\))?\s*)?(?:->\s*[^{}]+\s*)?\{"
)


def _parse_cpp(document: ExtractedDocument) -> Iterator[DocumentBlock]:
    builder = _BlockBuilder(document.path)
    source = _source_text(document)
    lines = source.splitlines(keepends=True)
    language = CPP_LANGUAGE[document.file_type]
    containers: list[tuple[int, int, str, str]] = []
    candidates: list[tuple[int, int, str, tuple[str, ...]]] = []
    covered: set[int] = set()

    try:
        for position, raw_line in enumerate(lines):
            stripped = raw_line.strip()
            if stripped.startswith("#include"):
                candidates.append((position, position, "include", ()))
                covered.add(position)
                continue
            if stripped.startswith("#define"):
                end = position
                while end + 1 < len(lines) and lines[end].rstrip().endswith("\\"):
                    end += 1
                candidates.append((position, end, "macro", ()))
                covered.update(range(position, end + 1))
                continue
            container_match = _CPP_CONTAINER.search(stripped)
            if container_match:
                kind, name = container_match.groups()
                end = _closing_line(lines, position)
                containers.append((position, end, kind, name))
                block_end = end if kind == "enum" else position
                candidates.append((position, block_end, kind, (name,)))
                covered.update(range(position, block_end + 1))

        for position, raw_line in enumerate(lines):
            match = _CPP_FUNCTION.search(raw_line.strip())
            if not match or match.group(1) in {"if", "for", "while", "switch", "catch"}:
                continue
            name_parts = tuple(part for part in match.group(1).split("::") if part)
            parents = tuple(
                name
                for start, end, kind, name in containers
                if kind in {"namespace", "class", "struct"} and start < position <= end
            )
            symbols = list(parents)
            for name in name_parts:
                if not symbols or symbols[-1] != name:
                    symbols.append(name)
            end = _closing_line(lines, position)
            candidates.append(
                (position, end, "method" if parents else "function", tuple(symbols))
            )
            covered.update(range(position, end + 1))

        depth = 0
        for position, raw_line in enumerate(lines):
            stripped = raw_line.strip()
            if (
                depth == 0
                and position not in covered
                and stripped.endswith(";")
                and "(" not in stripped
                and not stripped.startswith(("using ", "typedef ", "static_assert"))
            ):
                name_match = re.search(r"([A-Za-z_]\w*)\s*(?:=[^;]*)?;$", stripped)
                symbol = (name_match.group(1),) if name_match else ()
                candidates.append((position, position, "global-variable", symbol))
                covered.add(position)
            depth += _brace_change(raw_line)
    except Exception:
        logger.exception("C/C++ 结构解析失败，按行降级：%s", document.path)
        yield from _fallback_lines(
            document,
            source,
            language=language,
            builder=builder,
            parser="fallback-line",
        )
        return

    run_start: int | None = None
    for position in range(len(lines) + 1):
        unknown = position < len(lines) and position not in covered and lines[position].strip()
        if unknown and run_start is None:
            run_start = position
        if run_start is not None and (not unknown or position - run_start >= 80):
            candidates.append((run_start, position - 1, "fallback-line", ()))
            run_start = position if unknown else None

    emitted: set[tuple[int, int, str, tuple[str, ...]]] = set()
    for start, end, kind, symbols in sorted(candidates):
        identity = (start, end, kind, symbols)
        if identity in emitted:
            continue
        emitted.add(identity)
        content = "".join(lines[start : end + 1]).rstrip()
        if content:
            yield builder.make(
                kind,
                content,
                language=language,
                symbol_path=symbols,
                start_line=start + 1,
                end_line=end + 1,
                hard_boundary_before=True,
                hard_boundary_after=True,
                parser="cpp-brace",
            )


def build_embedding_content(block: DocumentBlock, canonical: str | None = None) -> str:
    """Build the stable natural-text document input used by Qwen3 Embedding."""

    return build_document_embedding_input(
        content=block.content if canonical is None else canonical,
        path=block.path,
        block_type=block.block_type,
        language=block.language,
        heading_path=block.heading_path,
        symbol_path=block.symbol_path,
        start_line=block.start_line,
        end_line=block.end_line,
        page_number=block.page_number,
        record_path=block.record_path,
        slide_number=block.slide_number,
    )


def iter_document_blocks(
    document: ExtractedDocument,
    *,
    json_profile: JsonProfile | None = None,
    max_json_size: int = DEFAULT_MAX_JSON_SIZE,
    json_record_probe_size: int = DEFAULT_JSON_RECORD_PROBE_SIZE,
) -> Iterator[DocumentBlock]:
    """Parse one extracted document without executing source code."""

    parsers = {
        "md": _parse_markdown,
        "txt": _parse_text,
        "pdf": _parse_pdf,
        "pptx": _parse_pptx,
        "py": _parse_python,
    }
    if document.file_type == "json":
        if json_profile is None:
            raise ValueError("索引 JSON 文件必须提供 --json-config 配置文件")
        yield from _parse_json(
            document,
            json_profile,
            max_json_size=max_json_size,
            record_probe_size=json_record_probe_size,
        )
        return
    if document.file_type in CPP_LANGUAGE:
        yield from _parse_cpp(document)
        return
    parser = parsers.get(document.file_type)
    if parser is None:
        raise ValueError(f"不支持的文档类型：{document.file_type}")
    yield from parser(document)
