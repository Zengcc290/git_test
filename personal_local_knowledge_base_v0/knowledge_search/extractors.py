"""TXT、Markdown、PDF 和 PPTX 的流式文本抽取。"""

from __future__ import annotations

# codecs 提供支持跨块解码的增量解码器。
import codecs
# hashlib 用于计算源文件指纹，实现增量索引。
import hashlib
# logging 负责记录编码降级、PDF 页失败等可诊断信息。
import logging
# Path 提供跨平台的文件系统访问。
from pathlib import Path
from collections.abc import Iterator

# 抽取结果统一使用这个数据模型传递给清洗和索引模块。
from .json_parser import (
    DEFAULT_MAX_JSON_SIZE,
    DEFAULT_JSON_RECORD_PROBE_SIZE,
    JsonProfile,
    iter_json_text,
)
from .models import ExtractedDocument


# 目前支持纯文本、Markdown、带文本层的 PDF、PPTX 和配置驱动的 JSON。
SUPPORTED_SUFFIXES = {".txt", ".md", ".pdf", ".pptx", ".json"}
# 单次读取的字节数；它限制 TXT/Markdown 的读取缓存大小。
DEFAULT_READ_SIZE = 64 * 1024
# 为当前模块创建独立 logger，便于按模块筛选日志。
logger = logging.getLogger(__name__)


def _choose_text_encoding(first_bytes: bytes, path: Path) -> str:
    """根据文件开头选择文本编码，避免读取完整文件才能判断编码。"""

    # UTF-8-sig 可以同时处理普通 UTF-8 和带 BOM 的 UTF-8 文件。
    utf8_decoder = codecs.getincrementaldecoder("utf-8-sig")(errors="strict")
    try:
        # final=False 允许首个块末尾正好截断一个多字节字符。
        utf8_decoder.decode(first_bytes, final=False)
        return "utf-8-sig"
    except UnicodeDecodeError:
        # UTF-8 失败时尝试中文 Windows 文档常见的 GB18030。
        logger.warning("文件 %s 不是 UTF-8，尝试使用 GB18030 解码", path)

    gb_decoder = codecs.getincrementaldecoder("gb18030")(errors="strict")
    try:
        gb_decoder.decode(first_bytes, final=False)
        return "gb18030"
    except UnicodeDecodeError:
        # 两种编码都无法确认时使用 replacement，保证索引流程不中断。
        logger.warning("文件 %s 无法可靠解码，将使用 UTF-8 replacement", path)
        return "utf-8"


def _iter_text_file(path: Path, read_size: int = DEFAULT_READ_SIZE) -> Iterator[str]:
    """以增量解码方式读取 TXT/Markdown，不把整个文件放进内存。"""

    if read_size <= 0:
        raise ValueError("read_size 必须大于 0")

    # 先读取一个固定大小的块用于编码判断，而不是调用 read_bytes。
    with path.open("rb") as stream:
        first_bytes = stream.read(read_size)
        if not first_bytes:
            return

        encoding = _choose_text_encoding(first_bytes, path)
        # errors=replace 是最后一道保护，避免后续某个坏字节让整次索引失败。
        decoder = codecs.getincrementaldecoder(encoding)(errors="replace")

        # 先输出已读取的第一块，再循环读取后续固定大小的块。
        first_text = decoder.decode(first_bytes, final=False)
        if first_text:
            yield first_text

        while True:
            data = stream.read(read_size)
            if not data:
                break
            text = decoder.decode(data, final=False)
            if text:
                yield text

        # final=True 刷出可能暂存在解码器中的半个字符。
        tail = decoder.decode(b"", final=True)
        if tail:
            yield tail


def _iter_pdf_text(path: Path) -> Iterator[str]:
    """逐页抽取 PDF 文本，每次只把当前页交给下游。"""

    try:
        # 延迟导入，保证只使用 TXT/Markdown/PPTX 时不受 PDF 依赖影响。
        from pypdf import PdfReader
    except ImportError as exc:  # pragma: no cover - requirements 会安装 pypdf
        raise RuntimeError("PDF 抽取需要 pypdf，请先安装 requirements.txt") from exc

    # PdfReader 负责解析 PDF 文件结构和页面对象。
    reader = PdfReader(str(path))
    for page_number, page in enumerate(reader.pages, start=1):
        try:
            # 某些页面没有文本层，extract_text 会返回 None。
            page_text = page.extract_text() or ""
            if page_text:
                # 一页一页交给下游，避免把所有页面拼成一个大字符串。
                yield page_text
                yield "\n\n"
        except Exception:
            # 单页失败时记录堆栈并继续读取其余页面。
            logger.exception("抽取 PDF %s 的第 %s 页失败", path, page_number)


def _extract_pptx_shape_text(shape, group_shape_type) -> list[str]:
    """递归抽取一个 PPTX 图形中的文本框、表格和组合图形文字。"""

    # 这里的列表只保存一张幻灯片内一个图形的文字，大小通常远小于整个演示文稿。
    texts: list[str] = []

    # 文本框、标题、占位符等对象通常通过 has_text_frame 暴露文字。
    if getattr(shape, "has_text_frame", False):
        text = shape.text.strip()
        if text:
            texts.append(text)

    # 表格对象没有普通文本框，需要逐行读取每个单元格的文字。
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                # 用制表符连接同一行，既保留列边界又方便全文搜索。
                texts.append("\t".join(cells))

    # 组合图形自身可能没有文本，但其子图形可能包含文本，需要递归处理。
    if getattr(shape, "shape_type", None) == group_shape_type:
        for child in shape.shapes:
            texts.extend(_extract_pptx_shape_text(child, group_shape_type))

    return texts


def _iter_pptx_text(path: Path) -> Iterator[str]:
    """逐张幻灯片抽取 PPTX 文本框、标题、表格和组合图形文字。"""

    try:
        # 延迟导入，让只处理 TXT/Markdown/PDF 时不受 PPTX 依赖影响。
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError as exc:  # pragma: no cover - requirements 会安装 python-pptx
        raise RuntimeError("PPTX 抽取需要 python-pptx，请先安装 requirements.txt") from exc

    # Presentation 负责读取 PPTX 压缩包内的演示文稿结构。
    presentation = Presentation(str(path))
    for slide_number, slide in enumerate(presentation.slides, start=1):
        # 只暂存当前幻灯片内容，处理完马上交给下游分段器。
        slide_texts: list[str] = []
        for shape in slide.shapes:
            # 按页面中的图形顺序读取，尽量保持用户看到的阅读顺序。
            slide_texts.extend(_extract_pptx_shape_text(shape, MSO_SHAPE_TYPE.GROUP))
        if slide_texts:
            # 加入页码标记，搜索结果中可以快速定位命中所在幻灯片。
            yield f"[Slide {slide_number}]\n" + "\n".join(slide_texts)
            # 用空行隔开幻灯片，避免上一页末尾和下一页开头发生文本粘连。
            yield "\n\n"


def extract_document(
    path: Path,
    *,
    parser_fingerprint: str = "",
) -> ExtractedDocument:
    """读取文件元数据并流式计算哈希，不在这里一次性抽取全文。"""

    # 展开用户目录符号并转成绝对路径，确保数据库键稳定。
    path = path.expanduser().resolve()
    if not path.is_file():
        raise FileNotFoundError(f"文件不存在：{path}")
    # 统一使用小写后缀，兼容 .TXT、.Md 等大小写变体。
    suffix = path.suffix.lower()
    if suffix not in SUPPORTED_SUFFIXES:
        raise ValueError(f"不支持的文件类型：{path.suffix or '<无扩展名>'}")

    # 使用固定大小的块计算哈希，避免大文件的 bytes 对象瞬间占满内存。
    hasher = hashlib.sha256()
    with path.open("rb") as stream:
        while True:
            data = stream.read(DEFAULT_READ_SIZE)
            if not data:
                break
            hasher.update(data)

    # stat 只读取文件元数据；text 保留为 None，正文由 iter_document_text 延迟读取。
    stat = path.stat()
    return ExtractedDocument(
        path=path,
        file_type=suffix[1:],
        text=None,
        sha256=hasher.hexdigest(),
        size=stat.st_size,
        modified_ns=stat.st_mtime_ns,
        parser_fingerprint=parser_fingerprint,
    )


def iter_document_text(
    document: ExtractedDocument,
    *,
    read_size: int = DEFAULT_READ_SIZE,
    json_profile: JsonProfile | None = None,
    max_json_size: int = DEFAULT_MAX_JSON_SIZE,
    json_record_probe_size: int = DEFAULT_JSON_RECORD_PROBE_SIZE,
) -> Iterator[str]:
    """根据文档类型按块、按页或按幻灯片产生文本。"""

    # 手工构造的测试对象可以直接提供 text，兼容旧的调用方式。
    if document.text is not None:
        yield document.text
        return

    # 普通文本文件走固定大小的增量解码器。
    if document.file_type in {"txt", "md"}:
        yield from _iter_text_file(document.path, read_size)
        return
    # PDF 按页产生文本，PPTX 按幻灯片产生文本。
    if document.file_type == "pdf":
        yield from _iter_pdf_text(document.path)
        return
    if document.file_type == "pptx":
        yield from _iter_pptx_text(document.path)
        return
    if document.file_type == "json":
        if json_profile is None:
            raise ValueError("索引 JSON 文件必须提供 --json-config 配置文件")
        yield from iter_json_text(
            document.path,
            json_profile,
            max_size=max_json_size,
            record_probe_size=json_record_probe_size,
        )
        return

    # extract_document 已经做过后缀校验，这里是对手工对象的额外保护。
    raise ValueError(f"不支持的文档类型：{document.file_type}")
