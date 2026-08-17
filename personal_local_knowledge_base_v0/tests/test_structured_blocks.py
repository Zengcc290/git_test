import importlib.util
import json
import sqlite3
import tempfile
import unittest
from pathlib import Path

from knowledge_search.block_parsing import iter_document_blocks
from knowledge_search.chunking import iter_chunk_blocks
from knowledge_search.database import KnowledgeBase
from knowledge_search.extractors import extract_document
from knowledge_search.indexer import index_paths
from knowledge_search.json_parser import JsonField, JsonProfile
from knowledge_search.rag.retriever import KeywordRetriever


HAS_PYPDF = importlib.util.find_spec("pypdf") is not None
HAS_PPTX = importlib.util.find_spec("pptx") is not None


def _profile() -> JsonProfile:
    return JsonProfile(
        name="users",
        record_path="$.users[*]",
        index_mode="record",
        fields=(
            JsonField(path="id", name="id"),
            JsonField(path="name", name="name"),
            JsonField(path="profile.city", name="profile.city"),
        ),
        separator="\n",
        filters=(),
        fingerprint="users-profile",
    )


class StructuredBlockTests(unittest.TestCase):
    def test_markdown_heading_example_and_node_types(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            path = Path(temp_dir) / "database.md"
            path.write_text(
                "# 数据库设计\n\n"
                "## documents 表\n\n这里是文档说明。\n\n"
                "- id\n- path\n\n"
                "> 引用说明\n\n"
                "```sql\nSELECT * FROM documents;\n```\n\n"
                "| 字段 | 类型 |\n| --- | --- |\n| id | INTEGER |\n\n"
                "## chunks 表\n\n这里是分段说明。\n",
                encoding="utf-8",
            )

            blocks = list(iter_document_blocks(extract_document(path)))

        self.assertEqual(blocks[0].block_type, "heading")
        self.assertEqual(blocks[0].heading_path, ("数据库设计",))
        self.assertEqual(blocks[0].content, "")
        self.assertEqual(blocks[1].heading_path, ("数据库设计", "documents 表"))
        self.assertEqual(blocks[1].content, "这里是文档说明。")
        self.assertTrue({"paragraph", "list", "quote", "code", "table"}.issubset(
            {block.block_type for block in blocks}
        ))
        self.assertEqual(blocks[-1].heading_path, ("数据库设计", "chunks 表"))
        self.assertEqual(blocks[-1].content, "这里是分段说明。")

    def test_canonical_and_embedding_content_are_stored_separately(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            path = root / "note.md"
            path.write_text("# 数据库设计\n\n这里是原始正文。", encoding="utf-8")

            with KnowledgeBase(root / "knowledge.db") as knowledge_base:
                stats = index_paths(knowledge_base, [path])
                row = knowledge_base.connection.execute(
                    "SELECT canonical_content, embedding_content FROM chunks"
                ).fetchone()
                result = knowledge_base.search("数据库设计")[0]
                retrieval = KeywordRetriever(knowledge_base, top_k=1).retrieve(
                    "数据库设计"
                )

            self.assertEqual(stats.indexed, 1)
            self.assertEqual(row["canonical_content"], "这里是原始正文。")
            self.assertNotIn("文档主题", row["canonical_content"])
            self.assertIn("结构：数据库设计", row["embedding_content"])
            self.assertIn("原文内容：", row["embedding_content"])
            self.assertEqual(result.content, "这里是原始正文。")
            self.assertIn("标题：数据库设计", retrieval.context)
            self.assertNotIn("文档主题：", retrieval.chunks[0].content)

    def test_json_records_keep_paths_and_hard_windows(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            path = root / "users.json"
            path.write_text(
                json.dumps(
                    {
                        "users": [
                            {"id": 1, "name": "Alice", "profile": {"city": "Beijing"}},
                            {"id": 2, "name": "Bob", "profile": {"city": "Shanghai"}},
                        ]
                    },
                    ensure_ascii=False,
                ),
                encoding="utf-8",
            )
            profile = _profile()
            document = extract_document(path, parser_fingerprint=profile.fingerprint)
            blocks = list(iter_document_blocks(document, json_profile=profile))

            with KnowledgeBase(root / "knowledge.db") as knowledge_base:
                index_paths(knowledge_base, [path], json_profile=profile)
                alice = knowledge_base.search("Alice")[0]
                window = knowledge_base.chunk_window(alice.chunk_id, radius=4)

        self.assertEqual([block.record_path for block in blocks], ["users[0]", "users[1]"])
        self.assertIn("id: 1", blocks[0].content)
        self.assertIn("profile.city: Beijing", blocks[0].content)
        self.assertTrue(all(block.hard_boundary_before for block in blocks))
        self.assertTrue(all(block.hard_boundary_after for block in blocks))
        self.assertEqual(alice.record_path, "users[0]")
        self.assertEqual([chunk.record_path for chunk in window], ["users[0]"])

    def test_python_ast_symbols_source_metadata_and_fallback(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            valid = root / "database.py"
            valid.write_text(
                '"""Database module."""\n'
                "import os\n\n"
                "CONSTANT = 1\n\n"
                "class KnowledgeBase:\n"
                "    def replace_document(self, path: str, force=False):\n"
                '        """Atomically replace one document."""\n'
                "        # Keep the source unchanged.\n"
                "        return path\n\n"
                "def index_paths(paths):\n"
                "    return list(paths)\n",
                encoding="utf-8",
            )
            invalid = root / "broken.py"
            invalid.write_text("def broken(:\n    pass\n", encoding="utf-8")

            valid_document = extract_document(valid)
            valid_blocks = list(iter_document_blocks(valid_document))
            method = next(
                block
                for block in valid_blocks
                if block.symbol_path == ("KnowledgeBase", "replace_document")
            )
            invalid_document = extract_document(invalid)
            invalid_blocks = list(iter_document_blocks(invalid_document))

            with KnowledgeBase(root / "knowledge.db") as knowledge_base:
                stats = index_paths(knowledge_base, [valid, invalid])
                parsers = {
                    document.filename: document.parser
                    for document in knowledge_base.list_documents()
                }

        self.assertEqual(valid_document.parser, "python-ast")
        self.assertEqual(method.block_type, "method")
        self.assertEqual(method.parameters, ("self", "path", "force"))
        self.assertEqual(method.docstring, "Atomically replace one document.")
        self.assertIn("# Keep the source unchanged.", method.comments)
        self.assertIn("def replace_document", method.content)
        self.assertIsNotNone(method.start_line)
        self.assertEqual(invalid_document.parser, "fallback-line")
        self.assertTrue(invalid_blocks)
        self.assertTrue(all(block.parser == "fallback-line" for block in invalid_blocks))
        self.assertEqual(stats.indexed, 2)
        self.assertEqual(parsers["broken.py"], "fallback-line")

    def test_cpp_namespace_class_and_method_symbols(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            path = Path(temp_dir) / "indexer.cpp"
            path.write_text(
                "#include <vector>\n"
                "#define INDEX_VERSION 4\n\n"
                "namespace knowledge {\n"
                "class Indexer {\n"
                "public:\n"
                "    void run() {\n"
                "        int error_code = 0;\n"
                "    }\n"
                "};\n"
                "}\n",
                encoding="utf-8",
            )

            blocks = list(iter_document_blocks(extract_document(path)))

        method = next(block for block in blocks if block.block_type == "method")
        self.assertEqual(method.symbol_path, ("knowledge", "Indexer", "run"))
        self.assertEqual(method.start_line, 7)
        self.assertIn("void run()", method.content)
        self.assertIn("include", {block.block_type for block in blocks})
        self.assertIn("macro", {block.block_type for block in blocks})

    def test_large_structural_block_splits_without_losing_metadata(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            path = Path(temp_dir) / "large.md"
            path.write_text("# Topic\n\n" + "abcdefghij" * 20, encoding="utf-8")
            blocks = iter_document_blocks(extract_document(path))
            chunks = list(iter_chunk_blocks(blocks, chunk_size=50, overlap=5))

        self.assertGreater(len(chunks), 1)
        self.assertTrue(all(chunk.heading_path == ("Topic",) for chunk in chunks))
        self.assertTrue(all("结构：Topic" in chunk.embedding_content for chunk in chunks))
        self.assertTrue(all("结构：" not in chunk.canonical_content for chunk in chunks))

    @unittest.skipUnless(HAS_PYPDF, "pypdf 未安装")
    def test_pdf_blocks_keep_page_numbers_and_boundaries(self):
        from pypdf import PdfWriter
        from pypdf.generic import DecodedStreamObject, DictionaryObject, NameObject

        with tempfile.TemporaryDirectory() as temp_dir:
            path = Path(temp_dir) / "pages.pdf"
            writer = PdfWriter()
            for text in (b"First page text", b"Second page text"):
                page = writer.add_blank_page(width=612, height=792)
                font = DictionaryObject(
                    {
                        NameObject("/Type"): NameObject("/Font"),
                        NameObject("/Subtype"): NameObject("/Type1"),
                        NameObject("/BaseFont"): NameObject("/Helvetica"),
                    }
                )
                page[NameObject("/Resources")] = DictionaryObject(
                    {NameObject("/Font"): DictionaryObject({NameObject("/F1"): font})}
                )
                stream = DecodedStreamObject()
                stream.set_data(b"BT /F1 12 Tf 72 720 Td (" + text + b") Tj ET")
                page[NameObject("/Contents")] = stream
            with path.open("wb") as handle:
                writer.write(handle)

            blocks = list(iter_document_blocks(extract_document(path)))

        self.assertEqual([block.page_number for block in blocks], [1, 2])
        self.assertTrue(all(block.hard_boundary_before for block in blocks))
        self.assertTrue(all(block.hard_boundary_after for block in blocks))

    @unittest.skipUnless(HAS_PPTX, "python-pptx 未安装")
    def test_pptx_blocks_keep_slide_shape_title_text_and_table(self):
        from pptx import Presentation
        from pptx.util import Inches

        with tempfile.TemporaryDirectory() as temp_dir:
            path = Path(temp_dir) / "slides.pptx"
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[5])
            slide.shapes.title.text = "Structured slides"
            text_box = slide.shapes.add_textbox(
                Inches(1), Inches(1.5), Inches(6), Inches(1)
            )
            text_box.text = "Slide body"
            table = slide.shapes.add_table(
                1, 2, Inches(1), Inches(3), Inches(5), Inches(1)
            ).table
            table.cell(0, 0).text = "field"
            table.cell(0, 1).text = "value"
            presentation.save(path)

            blocks = list(iter_document_blocks(extract_document(path)))

        self.assertTrue(all(block.slide_number == 1 for block in blocks))
        self.assertTrue(all(block.shape_index is not None for block in blocks))
        self.assertIn("title", {block.block_type for block in blocks})
        self.assertIn("text-box", {block.block_type for block in blocks})
        self.assertIn("table", {block.block_type for block in blocks})
        self.assertTrue(blocks[0].hard_boundary_before)
        self.assertTrue(blocks[-1].hard_boundary_after)


class StructuredDatabaseMigrationTests(unittest.TestCase):
    def test_legacy_v0_database_is_migrated_without_replacing_original_fts(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            db_path = Path(temp_dir) / "legacy.db"
            connection = sqlite3.connect(db_path)
            connection.executescript(
                """
                CREATE TABLE documents (
                    id INTEGER PRIMARY KEY,
                    path TEXT NOT NULL UNIQUE,
                    filename TEXT NOT NULL,
                    file_type TEXT NOT NULL,
                    sha256 TEXT NOT NULL,
                    size INTEGER NOT NULL,
                    modified_ns INTEGER NOT NULL,
                    parser_fingerprint TEXT NOT NULL DEFAULT '',
                    indexed_at TEXT NOT NULL DEFAULT CURRENT_TIMESTAMP
                );
                CREATE TABLE chunks (
                    id INTEGER PRIMARY KEY,
                    document_id INTEGER NOT NULL REFERENCES documents(id) ON DELETE CASCADE,
                    chunk_index INTEGER NOT NULL,
                    content TEXT NOT NULL,
                    start_offset INTEGER NOT NULL DEFAULT 0,
                    UNIQUE(document_id, chunk_index)
                );
                CREATE VIRTUAL TABLE chunks_fts USING fts5(
                    content,
                    content='chunks',
                    content_rowid='id',
                    tokenize='unicode61'
                );
                INSERT INTO documents(
                    id, path, filename, file_type, sha256, size, modified_ns
                ) VALUES (1, 'legacy.md', 'legacy.md', 'md', 'old', 11, 1);
                INSERT INTO chunks(id, document_id, chunk_index, content, start_offset)
                VALUES (1, 1, 0, 'legacy searchable', 0);
                INSERT INTO chunks_fts(rowid, content) VALUES (1, 'legacy searchable');
                """
            )
            connection.commit()
            connection.close()

            with KnowledgeBase(db_path) as knowledge_base:
                row = knowledge_base.connection.execute(
                    """
                    SELECT canonical_content, embedding_content
                    FROM chunks WHERE id = 1
                    """
                ).fetchone()
                original_fts_sql = knowledge_base.connection.execute(
                    "SELECT sql FROM sqlite_master WHERE name = 'chunks_fts'"
                ).fetchone()[0]
                results = knowledge_base.search("legacy")

        self.assertEqual(row["canonical_content"], "legacy searchable")
        self.assertEqual(row["embedding_content"], "legacy searchable")
        self.assertIn("content='chunks'", original_fts_sql)
        self.assertEqual(results[0].content, "legacy searchable")


if __name__ == "__main__":
    unittest.main()
