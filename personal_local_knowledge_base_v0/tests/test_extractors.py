import importlib.util
import tempfile
import unittest
from pathlib import Path

# 被测函数：元数据抽取和按块读取正文的迭代器。
from knowledge_search.extractors import extract_document, iter_document_text

HAS_PYPDF = importlib.util.find_spec("pypdf") is not None
HAS_PPTX = importlib.util.find_spec("pptx") is not None


class ExtractorTests(unittest.TestCase):
    @unittest.skipUnless(HAS_PYPDF, "pypdf 未安装")
    def test_extracts_text_from_pdf_text_layer(self):
        # 这个测试只验证文本层抽取，不涉及 OCR。
        from pypdf import PdfWriter
        from pypdf.generic import DecodedStreamObject, DictionaryObject, NameObject

        with tempfile.TemporaryDirectory() as temp_dir:
            path = Path(temp_dir) / "note.pdf"
            # 创建一个空白页面，再手动附加 Helvetica 字体和内容流。
            writer = PdfWriter()
            page = writer.add_blank_page(width=612, height=792)
            font = DictionaryObject(
                {
                    NameObject("/Type"): NameObject("/Font"),
                    NameObject("/Subtype"): NameObject("/Type1"),
                    NameObject("/BaseFont"): NameObject("/Helvetica"),
                }
            )
            page[NameObject("/Resources")] = DictionaryObject(
                {NameObject("/Font"): DictionaryObject({NameObject("/F1"): font})}
            )
            stream = DecodedStreamObject()
            stream.set_data(b"BT /F1 12 Tf 72 720 Td (PDF search works) Tj ET")
            page[NameObject("/Contents")] = stream
            # 将构造好的 PDF 写入临时路径。
            with path.open("wb") as handle:
                writer.write(handle)

            # 抽取后应识别 PDF 类型；正文通过迭代器按块取得。
            document = extract_document(path)
            self.assertEqual(document.file_type, "pdf")
            self.assertIsNone(document.text)
            self.assertIn("PDF search works", "".join(iter_document_text(document)))

    @unittest.skipUnless(HAS_PPTX, "python-pptx 未安装")
    def test_extracts_text_and_table_from_pptx(self):
        # 生成一个包含文本框和表格的最小演示文稿，验证 PPTX 抽取路径。
        from pptx import Presentation
        from pptx.util import Inches

        with tempfile.TemporaryDirectory() as temp_dir:
            path = Path(temp_dir) / "note.pptx"
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])

            # 文本框模拟标题或正文内容。
            text_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
            text_box.text = "PPTX search works"

            # 表格内容应按行抽取，并参与全文搜索。
            table = slide.shapes.add_table(2, 2, Inches(1), Inches(2), Inches(5), Inches(2)).table
            table.cell(0, 0).text = "反应物"
            table.cell(0, 1).text = "氧化剂"
            table.cell(1, 0).text = "还原剂"
            table.cell(1, 1).text = "电子转移"
            presentation.save(path)

            document = extract_document(path)
            self.assertEqual(document.file_type, "pptx")
            self.assertIsNone(document.text)
            text = "".join(iter_document_text(document))
            self.assertIn("PPTX search works", text)
            self.assertIn("氧化剂", text)


if __name__ == "__main__":
    # 允许直接执行本测试文件。
    unittest.main()
